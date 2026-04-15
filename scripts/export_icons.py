"""
export_icons.py
アイコン・画像を個別PNGとして書き出す。
  - 埋め込み画像: python-pptxで直接バイナリ取得
  - external画像: フルスライドPNGからcropして取得
"""

import sys
import os
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from constants import emu_to_px, OUTPUT_WIDTH, OUTPUT_HEIGHT


def main(pptx_path, data_dir="data", assets_dir="assets"):
    raw_shapes_path = os.path.join(data_dir, "raw_shapes.json")
    if not os.path.exists(raw_shapes_path):
        print("[ERROR] raw_shapes.json が見つかりません", file=sys.stderr)
        sys.exit(1)

    with open(raw_shapes_path, encoding="utf-8") as f:
        raw_data = json.load(f)

    icons_dir = os.path.join(assets_dir, "icons")
    slides_dir = os.path.join(assets_dir, "slides")
    os.makedirs(icons_dir, exist_ok=True)

    prs = Presentation(pptx_path)
    oh = raw_data["output_height"]

    embedded_count = 0
    cropped_count = 0
    skipped_count = 0

    for slide_idx, slide in enumerate(prs.slides):
        slide_shapes = raw_data["slides"][slide_idx]["shapes"]

        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue

            # raw_shapes.json から対応するデータを取得
            shape_data = None
            for sd in slide_shapes:
                if sd["shape_id"] == shape.shape_id:
                    shape_data = sd
                    break

            if shape_data is None:
                continue

            top_pct = shape_data["y"] / oh * 100 if oh > 0 else 0

            # ストライプパターン画像を除外（top < 5% or top > 95%）
            if top_pct < 5 or top_pct > 95:
                skipped_count += 1
                continue

            name = shape.name.replace(" ", "_").replace("/", "_")
            filename = f"slide_{slide_idx+1:02d}_{name}.png"
            filepath = os.path.join(icons_dir, filename)

            # 埋め込み画像を試行
            try:
                blob = shape.image.blob
                if len(blob) == 0:
                    print(f"  [WARN] 0バイト画像をスキップ: {shape.name}")
                    skipped_count += 1
                    continue

                with open(filepath, "wb") as f:
                    f.write(blob)
                embedded_count += 1
                print(f"  [埋込] {filename} ({len(blob)} bytes)")

            except ValueError:
                # external/linked画像 → フルスライドPNGからcrop
                full_png_path = os.path.join(slides_dir, f"slide_{slide_idx+1:02d}_full.png")
                if not os.path.exists(full_png_path):
                    print(f"  [WARN] フルPNGが見つかりません: {full_png_path}")
                    # プレースホルダ情報をJSONに残す（座標は保持済み）
                    skipped_count += 1
                    continue

                try:
                    full_img = Image.open(full_png_path)
                    x = shape_data["x"]
                    y = shape_data["y"]
                    w = shape_data["w"]
                    h = shape_data["h"]
                    # クロップ範囲がはみ出さないようにクランプ
                    x2 = min(x + w, full_img.width)
                    y2 = min(y + h, full_img.height)
                    cropped = full_img.crop((max(0, x), max(0, y), x2, y2))
                    cropped.save(filepath)
                    cropped_count += 1
                    print(f"  [crop] {filename} ({w}x{h}px)")
                except Exception as e:
                    print(f"  [WARN] crop失敗: {shape.name}: {e}", file=sys.stderr)
                    skipped_count += 1

    # shape_dataにimage_filenameを更新してraw_shapes.jsonを上書き
    # (export先のパスをRemotionから参照するため)
    for slide in raw_data["slides"]:
        for sd in slide["shapes"]:
            if sd["is_picture"] and sd["image_content_type"] != "external":
                expected = os.path.join(icons_dir, sd["image_filename"] or "")
                if not os.path.exists(expected):
                    sd["image_filename"] = None

    with open(raw_shapes_path, "w", encoding="utf-8") as f:
        json.dump(raw_data, f, ensure_ascii=False, indent=2)

    print(f"\n完了: 埋込={embedded_count}, crop={cropped_count}, スキップ={skipped_count}")


if __name__ == "__main__":
    pptx_file = sys.argv[1] if len(sys.argv) > 1 else "input.pptx"
    main(pptx_file)
