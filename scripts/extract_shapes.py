"""
extract_shapes.py
PPTXから全シェイプの座標・属性を抽出し、raw_shapes.json に出力する。
グループシェイプは子要素を再帰展開する。
"""

import sys
import os
import json
import traceback
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# 同ディレクトリの constants を import
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from constants import emu_to_px, OUTPUT_WIDTH, OUTPUT_HEIGHT


def rgb_to_hex(rgb_color):
    """RGBColor を '#RRGGBB' 文字列に変換"""
    if rgb_color is None:
        return None
    try:
        return f"#{rgb_color.red:02X}{rgb_color.green:02X}{rgb_color.blue:02X}"
    except Exception:
        return None


def get_font_color_hex(font):
    """フォントの色を取得。テーマカラー参照の場合はRGB変換を試みる。"""
    try:
        if font.color and font.color.rgb:
            return rgb_to_hex(font.color.rgb)
    except AttributeError:
        pass
    try:
        if font.color and font.color.theme_color is not None:
            # テーマカラーの場合、type が RGB なら取れる
            if font.color.type is not None:
                return rgb_to_hex(font.color.rgb)
    except Exception:
        pass
    return None


def get_fill_color_hex(shape):
    """シェイプの塗りつぶし色を取得"""
    try:
        fill = shape.fill
        if fill.type is not None and fill.fore_color and fill.fore_color.rgb:
            return rgb_to_hex(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def get_line_color_hex(shape):
    """枠線の色を取得"""
    try:
        line = shape.line
        if line.color and line.color.rgb:
            return rgb_to_hex(line.color.rgb)
    except Exception:
        pass
    return None


def get_line_width_pt(shape):
    """枠線の幅をpt単位で取得"""
    try:
        if shape.line.width:
            return round(shape.line.width / 12700, 2)
    except Exception:
        pass
    return None


def extract_text_runs(shape):
    """テキストフレームからrun情報を抽出する"""
    if not hasattr(shape, "text_frame"):
        return None

    try:
        tf = shape.text_frame
    except Exception:
        return None

    runs = []
    paragraphs = list(tf.paragraphs)

    for p_idx, para in enumerate(paragraphs):
        for run in para.runs:
            font_size = None
            if run.font.size:
                font_size = round(run.font.size / 12700, 1)

            runs.append({
                "text": run.text,
                "font_size_pt": font_size,
                "font_color": get_font_color_hex(run.font),
                "bold": bool(run.font.bold),
                "italic": bool(run.font.italic),
                "font_name": run.font.name
            })

        # paragraph間の改行を追加（最後の段落以外）
        if p_idx < len(paragraphs) - 1:
            runs.append({
                "text": "\n",
                "font_size_pt": None,
                "font_color": None,
                "bold": False,
                "italic": False,
                "font_name": None
            })

    # フォントサイズがnullの場合、同シェイプ内で最も出現頻度が高いサイズを適用
    size_counts = {}
    for r in runs:
        if r["font_size_pt"] is not None:
            s = r["font_size_pt"]
            size_counts[s] = size_counts.get(s, 0) + 1

    if size_counts:
        most_common_size = max(size_counts, key=size_counts.get)
        for r in runs:
            if r["font_size_pt"] is None and r["text"] != "\n":
                r["font_size_pt"] = most_common_size

    return runs if runs else None


def extract_single_shape(shape, slide_index, sw, sh):
    """単一シェイプの情報を辞書として抽出する"""
    try:
        stype = str(shape.shape_type).split("(")[0].replace("MSO_SHAPE_TYPE.", "").strip()

        has_text = False
        text = None
        text_runs = None

        if hasattr(shape, "text_frame"):
            try:
                raw_text = shape.text_frame.text
                if raw_text and raw_text.strip():
                    has_text = True
                    text = raw_text
                    text_runs = extract_text_runs(shape)
            except Exception:
                pass

        is_picture = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        image_content_type = None
        image_filename = None

        if is_picture:
            try:
                image_content_type = shape.image.content_type
                ext = image_content_type.split("/")[-1].replace("jpeg", "jpg")
                image_filename = f"slide_{slide_index+1:02d}_{shape.name.replace(' ', '_')}.{ext}"
            except ValueError:
                # external/linked image
                image_content_type = "external"
                image_filename = f"slide_{slide_index+1:02d}_{shape.name.replace(' ', '_')}.png"

        return {
            "slide_index": slide_index,
            "shape_id": shape.shape_id,
            "name": shape.name,
            "shape_type": stype,
            "x": emu_to_px(shape.left, sw, OUTPUT_WIDTH),
            "y": emu_to_px(shape.top, sh, OUTPUT_HEIGHT),
            "w": emu_to_px(shape.width, sw, OUTPUT_WIDTH),
            "h": emu_to_px(shape.height, sh, OUTPUT_HEIGHT),
            "left_emu": shape.left,
            "top_emu": shape.top,
            "width_emu": shape.width,
            "height_emu": shape.height,
            "rotation": shape.rotation,
            "has_text": has_text,
            "text": text,
            "text_runs": text_runs,
            "fill_color": get_fill_color_hex(shape),
            "line_color": get_line_color_hex(shape),
            "line_width_pt": get_line_width_pt(shape),
            "is_picture": is_picture,
            "image_content_type": image_content_type,
            "image_filename": image_filename,
        }
    except Exception as e:
        print(f"  [WARN] シェイプ '{getattr(shape, 'name', '?')}' の抽出に失敗: {e}", file=sys.stderr)
        traceback.print_exc(file=sys.stderr)
        return None


def extract_shapes_recursive(shapes, slide_index, sw, sh, group_offset_left=0, group_offset_top=0):
    """シェイプを再帰的に抽出する。グループは展開して子要素を取得。"""
    result = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            g_left = shape.left
            g_top = shape.top
            for child in shape.shapes:
                child_data = extract_single_shape(child, slide_index, sw, sh)
                if child_data:
                    # グループのオフセットをEMU単位で加算してからpx変換
                    child_data["x"] = emu_to_px(child.left + g_left + group_offset_left, sw, OUTPUT_WIDTH)
                    child_data["y"] = emu_to_px(child.top + g_top + group_offset_top, sh, OUTPUT_HEIGHT)
                    child_data["left_emu"] = child.left + g_left + group_offset_left
                    child_data["top_emu"] = child.top + g_top + group_offset_top
                    result.append(child_data)
        else:
            data = extract_single_shape(shape, slide_index, sw, sh)
            if data:
                # 親グループからのオフセットを加算
                if group_offset_left or group_offset_top:
                    data["x"] = emu_to_px(shape.left + group_offset_left, sw, OUTPUT_WIDTH)
                    data["y"] = emu_to_px(shape.top + group_offset_top, sh, OUTPUT_HEIGHT)
                    data["left_emu"] = shape.left + group_offset_left
                    data["top_emu"] = shape.top + group_offset_top
                result.append(data)
    return result


def main(pptx_path, output_dir="data"):
    if not os.path.exists(pptx_path):
        print(f"[ERROR] ファイルが見つかりません: {pptx_path}", file=sys.stderr)
        sys.exit(1)

    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(pptx_path)
    sw = prs.slide_width
    sh = prs.slide_height

    slides_data = []
    for i, slide in enumerate(prs.slides):
        print(f"  スライド {i+1}/{len(prs.slides)} を処理中...")
        shapes = extract_shapes_recursive(slide.shapes, i, sw, sh)
        slides_data.append({
            "slide_index": i,
            "shapes": shapes,
        })

    output = {
        "source_file": os.path.basename(pptx_path),
        "slide_width_emu": sw,
        "slide_height_emu": sh,
        "output_width": OUTPUT_WIDTH,
        "output_height": OUTPUT_HEIGHT,
        "slides": slides_data,
    }

    output_path = os.path.join(output_dir, "raw_shapes.json")
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(output, f, ensure_ascii=False, indent=2)

    total_shapes = sum(len(s["shapes"]) for s in slides_data)
    print(f"\n完了: {len(slides_data)} スライド、{total_shapes} シェイプを抽出")
    print(f"出力: {output_path}")


if __name__ == "__main__":
    pptx_file = sys.argv[1] if len(sys.argv) > 1 else "input.pptx"
    out_dir = sys.argv[2] if len(sys.argv) > 2 else "data"
    main(pptx_file, out_dir)
